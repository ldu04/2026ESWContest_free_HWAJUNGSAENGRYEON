"""deck_audit.py — 덱 레이아웃 검증. **렌더된 글자 좌표**로만 판정한다.

왜 프레임 좌표로는 안 되나
--------------------------
pptx 의 텍스트 프레임 높이는 선언값이다. 자동확장된 글자는 프레임 밖으로 나가고,
프레임끼리 비교하면 그 이탈이 안 잡힌다. (발표5 본문 프레임은 H=0.05in 인데 실제 3줄을 그렸다.)
그래서 PDF 에서 **실제 글자 span** 을 읽어, 각 span 을 그린 텍스트 상자에 귀속시키고,
그 상자의 **렌더 사각형**(span 들의 합집합)으로만 판정한다.

두 가지를 본다
  ① 이탈  — 렌더 사각형이 **뒤에 깔린 배경 도형(카드/박스)** 밖으로 나갔는가
  ② 겹침  — 모든 사각형 쌍(렌더 텍스트 · 표 · 배경 도형)의 교차

    python scripts/deck_audit.py <pptx> <pdf> [--verbose]
"""
import sys, fitz
from pptx import Presentation
sys.stdout.reconfigure(encoding="utf-8")
I = lambda v: v/914400.0
PAD = 0.12          # 배경 도형 안쪽 최소 여백
EPS = 0.005         # 겹침으로 볼 최소 침범 깊이(in). 0.01in 도 화면에선 맞닿아 보인다.

def spans(page):
    out=[]
    for b in page.get_text("dict")["blocks"]:
        for l in b.get("lines",[]):
            for sp in l["spans"]:
                if sp["text"].strip():
                    x0,y0,x1,y1=sp["bbox"]
                    out.append([x0/72,y0/72,x1/72,y1/72,sp["text"]])
    return out

def table_drawn_rect(page, shape):
    """표의 **그려진** 사각형. 행 높이 합도 최소값일 뿐이라 실제 렌더와 다르다.
    표 테두리는 벡터로 그려지므로 PDF drawings 에서 그 가로선의 최하단을 읽는다."""
    L,T=I(shape.left),I(shape.top)
    R=L+sum(c.width for c in shape.table.columns)/914400.0
    B=T+sum(r.height for r in shape.table.rows)/914400.0
    lo=B
    for d in page.get_drawings():
        x0,y0,x1,y1=[v/72 for v in d["rect"]]
        if x1-x0 < (R-L)*0.5: continue            # 표 폭의 절반 이상인 가로 요소만
        if x0 < L-0.15 or x1 > R+0.15: continue
        if T-0.05 <= y0 <= T+(B-T)*3 + 2.0:
            lo=max(lo, y1)
    return L,T,R,lo

def rect_of(shape):
    """표는 선언 높이가 1.00 으로 고정돼 있어 쓸 수 없다 — 행 높이 합을 쓴다."""
    L,T=I(shape.left),I(shape.top)
    if shape.has_table:
        return L, T, L+sum(c.width for c in shape.table.columns)/914400.0, \
                     T+sum(r.height for r in shape.table.rows)/914400.0
    return L, T, L+I(shape.width), T+I(shape.height)

def inter(a,b):
    x=min(a[2],b[2])-max(a[0],b[0]); y=min(a[3],b[3])-max(a[1],b[1])
    return (x,y) if (x>0 and y>0) else None

def analyze(pptx, pdf):
    p=Presentation(pptx); doc=fitz.open(pdf)
    W,H=I(p.slide_width),I(p.slide_height)
    out=[]
    for n,s in enumerate(p.slides,1):
        sp=spans(doc[n-1])
        boxes=[t for t in s.shapes if t.has_text_frame and t.text_frame.text.strip()
               and None not in (t.left,t.top,t.width,t.height)]
        # 표 안의 글자는 표가 주인이다. 이걸 안 빼면 표 위 문단이 표 글자까지 삼켜
        # 렌더 사각형이 표 전체를 덮고, 있지도 않은 '문단↔표 겹침'이 잡힌다.
        tmap={id(x._element): table_drawn_rect(doc[n-1], x) for x in s.shapes if x.has_table}
        trects=list(tmap.values())
        def in_table(x):
            cx,cy=(x[0]+x[2])/2,(x[1]+x[3])/2
            return any(L-0.02<=cx<=R+0.02 and T-0.02<=cy<=B+0.02 for L,T,R,B in trects)
        # ── span → 텍스트 상자 귀속 ─────────────────────────────────────
        # 좌표만으로는 못 가른다. 상자가 자기 높이를 크게 넘겨 그려지면, 그 아래 상자가
        # 넘친 줄을 가로채고(=겹침이 사라짐), 반대로 위 상자가 아랫줄을 삼킨다(=없는 겹침).
        # 그래서 **글자 내용**으로 먼저 붙이고, 내용으로 못 가릴 때만 좌표로 붙인다.
        import re as _re
        norm=lambda z: _re.sub(r"\s+","", z)
        btxt={id(t._element): norm(t.text_frame.text) for t in boxes}
        rend={}
        for x in sp:
            if in_table(x): continue          # 표 글자는 표가 주인 — 위 문단이 삼키지 못하게
            key=norm(x[4])
            cand=[t for t in boxes
                  if I(t.left)-0.10<=x[0] and x[2]<=I(t.left+t.width)+0.35]
            if not cand: continue
            # 내용 일치는 **6자 이상**일 때만 신뢰한다. 짧은 조각('·', 숫자)은 여러 상자에
            # 다 들어 있어 오히려 엉뚱한 상자에 붙는다. 유일하게 일치할 때만 좌표를 무시한다.
            byname=[t for t in cand if key in btxt[id(t._element)]] if len(key)>=6 else []
            if len(byname)==1:
                t=byname[0]
            else:
                ywin=[t for t in (byname or cand)
                      if I(t.top)<=x[1]+0.35 and x[1]<=I(t.top)+max(I(t.height),2.5)]
                if not ywin: continue
                def score(t):
                    d=x[1]-I(t.top)
                    return d if d>=0 else -d*10
                t=min(ywin, key=score)
            r=rend.setdefault(id(t._element), [t, x[0],x[1],x[2],x[3]])
            r[1]=min(r[1],x[0]); r[2]=min(r[2],x[1]); r[3]=max(r[3],x[2]); r[4]=max(r[4],x[3])
        # 배경 도형: 글자가 없고 면적이 큰 것
        cards=[c for c in s.shapes if not c.has_table and None not in (c.left,c.top,c.width,c.height)
               and c.width>0.9*914400 and c.height>0.45*914400
               and not (c.has_text_frame and c.text_frame.text.strip())]
        order={id(sh._element):i for i,sh in enumerate(s.shapes)}
        # ── ① 배경 도형 밖으로 나간 글자 ──
        for _,(t,x0,y0,x1,y1) in rend.items():
            back=[c for c in cards
                  if rect_of(c)[0]-0.05<=I(t.left) and I(t.left)<=rect_of(c)[2]+0.05
                  and rect_of(c)[1]-0.05<=I(t.top) < rect_of(c)[3]
                  and order[id(c._element)]<order[id(t._element)]]
            if not back: continue
            c=min(back, key=lambda z: z.width*z.height)
            L,T,R,B=rect_of(c)
            g=(y0-T, B-y1, x0-L, R-x1)
            if min(g)<PAD:
                out.append((n,"이탈" if min(g)<0 else "여백부족", t.name, c.name,
                            "상%+.2f 하%+.2f 좌%+.2f 우%+.2f"%g,
                            t.text_frame.text.strip().replace("\n"," ")[:30]))
        # ── ② 모든 사각형 쌍 교차 ──
        items=[(t.name, (x0,y0,x1,y1), "글자", id(t._element),
                t.text_frame.text.strip().replace("\n"," ")[:22])
               for _,(t,x0,y0,x1,y1) in rend.items()]
        items+=[(x.name, tmap[id(x._element)], "표", id(x._element), "") for x in s.shapes if x.has_table]
        for a in range(len(items)):
            for b in range(a+1,len(items)):
                A,B=items[a],items[b]
                v=inter(A[1],B[1])
                if v and min(v)>EPS:
                    out.append((n,"겹침",A[0],B[0],"침범 %.2f x %.2f in"%v,
                                (A[4] or B[4])[:30]))
    return out

if __name__=="__main__":
    res=analyze(sys.argv[1], sys.argv[2])
    from collections import Counter
    c=Counter(r[1] for r in res)
    print("이탈 %d · 여백부족 %d · 겹침 %d  (합 %d건)"
          % (c.get("이탈",0), c.get("여백부족",0), c.get("겹침",0), len(res)))
    for r in sorted(res, key=lambda z:(z[0],z[1])):
        print("  p%-2d %-6s %-9s ↔ %-9s  %-32s %s" % (r[0]-1,r[1],r[2],r[3],r[4],r[5]))
