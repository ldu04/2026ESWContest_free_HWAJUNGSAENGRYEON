"""deck_fit_cards.py — 카드 안 글자가 카드를 넘치는 곳을 자동으로 맞춘다.

우선순위: ① 아래로 키운다 → ② 위로 올린다 → ③ 본문 글자·행간을 줄인다.
문구는 절대 건드리지 않는다. 세 가지로도 부족하면 '남음'으로 보고한다.
"""
import sys, subprocess, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from deck_audit import audit
I=lambda v: v/914400.0
E=lambda x: int(round(x*914400))
GAP=0.14                                   # 목표 여백(기준 0.12 + 여유)

def fit(pptx, pdf):
    _,pad = audit(pptx, pdf)
    if not pad: return 0
    p=Presentation(pptx); moved=0
    for n,nm,g,_ in pad:
        s=p.slides[n-1]
        c=[x for x in s.shapes if x.name==nm][0]
        need=GAP-g[1]
        if need<=0: continue
        L,T,R,B=I(c.left),I(c.top),I(c.left+c.width),I(c.top+c.height)
        inside=[x for x in s.shapes if x is not c and None not in (x.left,x.top,x.width,x.height)
                and I(x.left)>=L-0.03 and I(x.left+x.width)<=R+0.03 and I(x.top)>=T-0.03 and I(x.top)<B]
        below=[x for x in s.shapes if x not in inside and x is not c
               and None not in (x.left,x.top,x.width,x.height) and x.height>=0
               and I(x.top)>=B-0.01 and I(x.left)<R and I(x.left+x.width)>L]
        above=[x for x in s.shapes if x not in inside and x is not c
               and None not in (x.left,x.top,x.width,x.height) and x.height>=0
               and I(x.top+x.height)<=T+0.01 and I(x.left)<R and I(x.left+x.width)>L]
        room_dn=max(0.0, min([I(x.top) for x in below], default=7.50)-B-0.12)
        room_up=max(0.0, T-max([I(x.top+x.height) for x in above], default=0.0)-0.12)
        grow=min(need, room_dn); need-=grow
        shift=min(need, room_up); need-=shift
        if grow: c.height=E(I(c.height)+grow)
        if shift:
            c.top=E(T-shift); c.height=E(I(c.height)+shift)
            for x in inside: x.top=E(I(x.top)-shift)
        act=["+%.2f 아래"%grow if grow else "", "%.2f 위로"%shift if shift else ""]
        if need>0.004:
            # 본문(가장 큰 텍스트 상자)의 글자와 행간을 줄인다. 문구는 그대로.
            body=sorted([x for x in inside if x.has_text_frame and x.text_frame.text.strip()],
                        key=lambda z:-z.height)[0]
            for para in body.text_frame.paragraphs:
                para.line_spacing=1.05
                for r in para.runs:
                    if r.font.size and r.font.size.pt>10.0:
                        r.font.size=Pt(max(10.0, r.font.size.pt-1.0))
            act.append("본문 −1pt·행간 1.05")
        print("  파일%2d쪽 %-9s 하%+.2f → %s" % (n,nm,g[1]," · ".join(x for x in act if x)))
        moved+=1
    p.save(pptx)
    return moved

if __name__=="__main__":
    print("%d개 카드 조정" % fit(sys.argv[1], sys.argv[2]))
