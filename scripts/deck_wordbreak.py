"""deck_wordbreak.py — 어절 중간에서 줄이 바뀐 곳을 원문 대조로 찾는다.

렌더된 줄을 이어붙이면 원문(공백 제거)과 같아야 한다. 줄이 바뀐 지점이
**원문에서 공백이 아닌 자리**면 낱말이 반토막 난 것이다. 눈으로 세지 않는다.

    python scripts/deck_wordbreak.py <pptx> <pdf>
"""
import sys, re, fitz
from pptx import Presentation
sys.stdout.reconfigure(encoding="utf-8")
I=lambda v: v/914400.0
NB=lambda s: re.sub(r"\s+","",s)

def lines_of(page):
    out=[]
    for b in page.get_text("dict")["blocks"]:
        for l in b.get("lines",[]):
            t="".join(sp["text"] for sp in l["spans"])
            if t.strip():
                x0=min(sp["bbox"][0] for sp in l["spans"])/72
                y0=min(sp["bbox"][1] for sp in l["spans"])/72
                x1=max(sp["bbox"][2] for sp in l["spans"])/72
                out.append((y0,x0,x1,t))
    return sorted(out)

def check(pptx,pdf):
    p=Presentation(pptx); doc=fitz.open(pdf); hits=[]
    for n,s in enumerate(p.slides,1):
        ls=lines_of(doc[n-1])
        used=[False]*len(ls)
        for sh in s.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip(): continue
            if None in (sh.left,sh.width): continue
            L,R=I(sh.left),I(sh.left+sh.width)
            for para in sh.text_frame.paragraphs:
                src=para.text
                if len(NB(src))<8: continue
                # 이 문단을 이룬 줄들을 순서대로 찾는다
                rest=NB(src); seq=[]
                for i,(y,x0,x1,t) in enumerate(ls):
                    if used[i] or not (L-0.15<=x0 and x1<=R+0.40): continue
                    k=NB(t)
                    if k and rest.startswith(k):
                        seq.append(i); used[i]=True; rest=rest[len(k):]
                        if not rest: break
                if rest or len(seq)<2: continue
                pos=0
                for i in seq[:-1]:
                    pos+=len(NB(ls[i][3]))
                    # 원문에서 pos 번째 '비공백' 문자 뒤가 공백인가
                    cnt=0; idx=0
                    for idx,ch in enumerate(src):
                        if not ch.isspace():
                            cnt+=1
                            if cnt==pos: break
                    tail=src[idx+1:idx+2]
                    if tail and not tail.isspace():
                        hits.append((n, sh.name, src[max(0,idx-14):idx+1]+" / "+src[idx+1:idx+15]))
    return hits

if __name__=="__main__":
    h=check(sys.argv[1], sys.argv[2])
    print("어절 중간 줄바꿈 %d건"%len(h))
    for n,nm,frag in h: print("   p%-2d %-9s %s"%(n-1,nm,frag))
