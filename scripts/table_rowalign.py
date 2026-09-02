"""table_rowalign.py — 표의 **행 안 세로 정렬을 픽셀로** 검사한다.

왜 속성 검사를 쓰면 안 되나 (2026-09-03)
----------------------------------------
`cell.vertical_anchor` 를 훑어 「전부 ctr」이라 보고했는데 렌더에서는 「31 m」이 위로
붙어 있었다. 원인은 정렬 속성이 아니라 **셀에 남은 <a:br/>** 이었다 — 내용이 2줄이 되어
가운데 정렬의 첫 줄이 올라간다. 속성은 둘 다 ctr 이라 영원히 못 잡는다.

무엇을 재나
-----------
셀 단위로 **보이는 글자의 합집합 상자 중심 y** 를 재서 같은 행끼리 비교한다.
정상적으로 두 줄로 접힌 셀은 두 줄의 합집합 중심이 행 중심과 같으므로 걸리지 않는다.
꼬리 <a:br/> 처럼 **빈 줄이 붙은 셀만** 중심이 위로 밀려 걸린다.

    python scripts/table_rowalign.py <pptx> <pdf> [허용px]
"""
import sys, fitz
from pptx import Presentation
sys.stdout.reconfigure(encoding="utf-8")
DPI, SW = 200, 13.333333
IN = 914400


def main(pptx, pdf, tol=2.0):
    pr = Presentation(pptx); doc = fitz.open(pdf)
    bad = []
    for n, s in enumerate(pr.slides):
        tables = [x for x in s.shapes if x.has_table]
        if not tables:
            continue
        page = doc[n]; pm = page.get_pixmap(dpi=DPI); S = pm.width / SW
        spans = []
        for b in page.get_text("dict")["blocks"]:
            for l in b.get("lines", []):
                for sp in l["spans"]:
                    if sp["text"].strip():
                        x0, y0, x1, y1 = [v / 72 for v in sp["bbox"]]
                        spans.append((x0, y0, x1, y1, sp["text"].strip()))
        for ti, sh in enumerate(tables):
            t = sh.table
            L = sh.left / IN
            edges = [L]
            for c in t.columns:
                edges.append(edges[-1] + c.width / IN)
            top = sh.top / IN
            bot = top + sum(r.height for r in t.rows) / IN + 1.2   # 렌더가 더 길 수 있다
            inside = [p for p in spans if L - .05 <= p[0] and p[2] <= edges[-1] + .10
                      and top - .10 <= p[1] <= bot]
            if not inside:
                continue
            inside.sort(key=lambda p: (p[1] + p[3]) / 2)
            groups, cur = [], []
            for p in inside:
                cy = (p[1] + p[3]) / 2
                if cur and (cy - (cur[-1][1] + cur[-1][3]) / 2) * S > 26:
                    groups.append(cur); cur = []
                cur.append(p)
            if cur:
                groups.append(cur)
            for gi, g in enumerate(groups):
                cols = {}
                for p in g:
                    cx = (p[0] + p[2]) / 2
                    ci = max(0, sum(1 for e in edges[1:-1] if e <= cx))
                    cols.setdefault(ci, []).append(p)
                if len(cols) < 2:
                    continue
                ctr = {ci: (min(q[1] for q in v) + max(q[3] for q in v)) / 2 * S
                       for ci, v in cols.items()}
                d = max(ctr.values()) - min(ctr.values())
                if d > tol:
                    bad.append((n, ti, gi, d, {ci: (round(y, 1), cols[ci][0][4][:16])
                                               for ci, y in sorted(ctr.items())}))
    if bad:
        print(" 쪽 | 표 | 행 | 편차px | 열별 중심 y · 첫 글자")
        print("----+----+----+--------+" + "-" * 46)
        for n, ti, gi, d, info in bad:
            print(" %2d | %2d | %2d | %6.1f | %s" % (n, ti, gi, d,
                  "  ".join("c%d %.0f %r" % (c, y, t) for c, (y, t) in info.items())[:70]))
    print("\n행 안 세로 어긋남 %d건 (허용 %.1fpx · 셀 단위 글자 합집합 중심)" % (len(bad), tol))
    return len(bad)


if __name__ == "__main__":
    sys.exit(0 if main(sys.argv[1], sys.argv[2],
                       float(sys.argv[3]) if len(sys.argv) > 3 else 2.0) == 0 else 1)
